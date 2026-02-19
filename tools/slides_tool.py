"""Slides Tool - Pembuatan presentasi dengan ekspor HTML/PPTX dan konten dinamis."""

import json
import logging
import os
import time
from typing import Optional

logger = logging.getLogger(__name__)


class Slide:
    def __init__(self, title: str, content: str, layout: str = "title_content", notes: str = ""):
        self.title = title
        self.content = content
        self.layout = layout
        self.notes = notes
        self.images: list[str] = []

    def to_dict(self) -> dict:
        return {
            "title": self.title,
            "content": self.content,
            "layout": self.layout,
            "notes": self.notes,
            "images": self.images,
        }


class Presentation:
    def __init__(self, title: str, author: str = "", theme: str = "modern"):
        self.title = title
        self.author = author
        self.theme = theme
        self.slides: list[Slide] = []
        self.created_at = time.time()

    def add_slide(self, title: str, content: str, layout: str = "title_content", notes: str = "") -> Slide:
        slide = Slide(title=title, content=content, layout=layout, notes=notes)
        self.slides.append(slide)
        return slide

    def to_dict(self) -> dict:
        return {
            "title": self.title,
            "author": self.author,
            "theme": self.theme,
            "slides": [s.to_dict() for s in self.slides],
            "slide_count": len(self.slides),
        }


class SlidesTool:
    def __init__(self, output_dir: str = "data/presentations", default_template: str = "modern"):
        self.output_dir = output_dir
        self.default_template = default_template
        self.presentations: list[Presentation] = []
        os.makedirs(output_dir, exist_ok=True)

    async def execute(self, plan: dict) -> str:
        intent = plan.get("intent", "")
        return (
            f"Slides tool siap. Intent: {intent}. "
            f"Operasi: create_presentation, add_slide, export."
        )

    def create_presentation(self, title: str, author: str = "", theme: Optional[str] = None) -> Presentation:
        pres = Presentation(title=title, author=author, theme=theme or self.default_template)
        self.presentations.append(pres)
        logger.info(f"Presentasi dibuat: '{title}'")
        return pres

    def add_slide_to_presentation(self, presentation: Presentation, title: str, content: str,
                                   layout: str = "title_content") -> Slide:
        slide = presentation.add_slide(title=title, content=content, layout=layout)
        logger.info(f"Slide ditambahkan: '{title}' ke presentasi '{presentation.title}'")
        return slide

    def generate_outline(self, topic: str, num_slides: int = 10) -> list[dict]:
        outline = [
            {"title": f"Judul: {topic}", "layout": "title"},
            {"title": "Pendahuluan", "layout": "title_content"},
            {"title": "Latar Belakang", "layout": "title_content"},
        ]
        for i in range(num_slides - 5):
            outline.append({"title": f"Poin {i + 1}", "layout": "title_content"})
        outline.append({"title": "Kesimpulan", "layout": "title_content"})
        outline.append({"title": "Terima Kasih", "layout": "title"})
        return outline

    def export_presentation(self, presentation: Presentation, filename: Optional[str] = None) -> str:
        if not filename:
            safe_title = presentation.title.replace(" ", "_").lower()
            filename = f"{safe_title}_{int(time.time())}.json"
        output_path = os.path.join(self.output_dir, filename)
        import json
        with open(output_path, "w") as f:
            json.dump(presentation.to_dict(), f, indent=2, ensure_ascii=False)
        logger.info(f"Presentasi diekspor: {output_path}")
        return output_path

    def add_slide(self, presentation, title: str, content: str, layout: str = "title_content") -> Slide:
        return self.add_slide_to_presentation(presentation, title, content, layout)

    def list_presentations(self) -> list[dict]:
        return [{"title": p.title, "slides": len(p.slides), "theme": p.theme} for p in self.presentations]

    def export_html(self, title_or_pres, output_path: Optional[str] = None) -> str:
        pres = None
        if isinstance(title_or_pres, Presentation):
            pres = title_or_pres
        elif isinstance(title_or_pres, str):
            for p in self.presentations:
                if p.title == title_or_pres:
                    pres = p
                    break
        if not pres:
            if self.presentations:
                pres = self.presentations[-1]
            else:
                return "Tidak ada presentasi untuk di-export."

        theme_colors = {
            "modern": {"bg": "#1a1a2e", "text": "#e4e4ef", "accent": "#6c5ce7", "slide_bg": "#16213e"},
            "dark": {"bg": "#0d1117", "text": "#c9d1d9", "accent": "#58a6ff", "slide_bg": "#161b22"},
            "light": {"bg": "#ffffff", "text": "#333333", "accent": "#0066cc", "slide_bg": "#f8f9fa"},
        }
        colors = theme_colors.get(pres.theme, theme_colors["modern"])

        slides_html = ""
        for i, slide in enumerate(pres.slides):
            content_html = slide.content.replace("\n", "<br>") if slide.content else ""
            if slide.layout == "title":
                slide_body = f'<div class="slide-center"><h1>{slide.title}</h1></div>'
            else:
                slide_body = f'<h2>{slide.title}</h2><div class="slide-body">{content_html}</div>'
            slides_html += f'''
    <div class="slide" id="slide-{i}">
      <div class="slide-number">{i+1}/{len(pres.slides)}</div>
      {slide_body}
    </div>'''

        html = f'''<!DOCTYPE html>
<html lang="id">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{pres.title}</title>
<style>
* {{ margin: 0; padding: 0; box-sizing: border-box; }}
body {{ background: {colors["bg"]}; color: {colors["text"]}; font-family: 'Segoe UI', system-ui, sans-serif; overflow: hidden; }}
.slide {{ display: none; width: 100vw; height: 100vh; padding: 60px 80px; position: relative; background: {colors["slide_bg"]}; }}
.slide.active {{ display: flex; flex-direction: column; justify-content: center; }}
.slide-center {{ display: flex; flex-direction: column; align-items: center; justify-content: center; height: 100%; text-align: center; }}
.slide h1 {{ font-size: 3em; color: {colors["accent"]}; margin-bottom: 0.5em; }}
.slide h2 {{ font-size: 2.2em; color: {colors["accent"]}; margin-bottom: 1em; border-bottom: 3px solid {colors["accent"]}; padding-bottom: 0.3em; }}
.slide-body {{ font-size: 1.4em; line-height: 1.8; }}
.slide-number {{ position: absolute; bottom: 20px; right: 30px; font-size: 0.9em; opacity: 0.5; }}
.controls {{ position: fixed; bottom: 20px; left: 50%; transform: translateX(-50%); display: flex; gap: 10px; z-index: 100; }}
.controls button {{ padding: 10px 24px; font-size: 1em; border: none; border-radius: 8px; cursor: pointer; background: {colors["accent"]}; color: white; }}
.controls button:hover {{ opacity: 0.85; }}
</style>
</head>
<body>
{slides_html}
<div class="controls">
  <button onclick="prev()">&#9664; Prev</button>
  <button onclick="next()">Next &#9654;</button>
</div>
<script>
let current = 0;
const slides = document.querySelectorAll('.slide');
function show(n) {{ slides.forEach(s => s.classList.remove('active')); if(slides[n]) slides[n].classList.add('active'); }}
function next() {{ if(current < slides.length-1) {{ current++; show(current); }} }}
function prev() {{ if(current > 0) {{ current--; show(current); }} }}
document.addEventListener('keydown', e => {{ if(e.key==='ArrowRight'||e.key===' ') next(); if(e.key==='ArrowLeft') prev(); }});
show(0);
</script>
</body>
</html>'''

        if not output_path:
            safe_title = pres.title.replace(" ", "_").lower()
            output_path = os.path.join(self.output_dir, f"{safe_title}_{int(time.time())}.html")

        with open(output_path, "w", encoding="utf-8") as f:
            f.write(html)

        logger.info(f"Presentasi HTML diekspor: {output_path}")
        return f"Presentasi '{pres.title}' berhasil di-export ke HTML: {output_path}"

    def export_pptx(self, title_or_pres, output_path: Optional[str] = None) -> str:
        pres = self._resolve_presentation(title_or_pres)
        if not pres:
            return "Tidak ada presentasi untuk di-export."

        try:
            from pptx import Presentation as PptxPresentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            return "python-pptx tidak terinstal. Jalankan: pip install python-pptx"

        pptx_pres = PptxPresentation()
        pptx_pres.slide_width = Inches(13.333)
        pptx_pres.slide_height = Inches(7.5)

        for slide in pres.slides:
            if slide.layout == "title":
                layout = pptx_pres.slide_layouts[0]
                pptx_slide = pptx_pres.slides.add_slide(layout)
                pptx_slide.shapes.title.text = slide.title
                if len(pptx_slide.placeholders) > 1:
                    pptx_slide.placeholders[1].text = slide.content or ""
            else:
                layout = pptx_pres.slide_layouts[1]
                pptx_slide = pptx_pres.slides.add_slide(layout)
                pptx_slide.shapes.title.text = slide.title
                if len(pptx_slide.placeholders) > 1:
                    tf = pptx_slide.placeholders[1].text_frame
                    tf.text = ""
                    for line in (slide.content or "").split("\n"):
                        p = tf.add_paragraph()
                        p.text = line
                        p.font.size = Pt(18)

            if slide.notes:
                notes_slide = pptx_slide.notes_slide
                notes_slide.notes_text_frame.text = slide.notes

        if not output_path:
            safe_title = pres.title.replace(" ", "_").lower()
            output_path = os.path.join(self.output_dir, f"{safe_title}_{int(time.time())}.pptx")

        pptx_pres.save(output_path)
        logger.info(f"Presentasi PPTX diekspor: {output_path}")
        return f"Presentasi '{pres.title}' berhasil di-export ke PPTX: {output_path}"

    def _resolve_presentation(self, title_or_pres) -> Optional[Presentation]:
        if isinstance(title_or_pres, Presentation):
            return title_or_pres
        elif isinstance(title_or_pres, str):
            for p in self.presentations:
                if p.title == title_or_pres:
                    return p
        if self.presentations:
            return self.presentations[-1]
        return None

    def update_slide(self, presentation: Presentation, index: int,
                     title: Optional[str] = None, content: Optional[str] = None,
                     layout: Optional[str] = None, notes: Optional[str] = None) -> dict:
        if index < 0 or index >= len(presentation.slides):
            return {"success": False, "error": f"Index {index} di luar jangkauan"}

        slide = presentation.slides[index]
        if title is not None:
            slide.title = title
        if content is not None:
            slide.content = content
        if layout is not None:
            slide.layout = layout
        if notes is not None:
            slide.notes = notes

        return {"success": True, "slide": slide.to_dict()}

    def remove_slide(self, presentation: Presentation, index: int) -> dict:
        if index < 0 or index >= len(presentation.slides):
            return {"success": False, "error": f"Index {index} di luar jangkauan"}
        removed = presentation.slides.pop(index)
        return {"success": True, "removed": removed.to_dict(), "remaining": len(presentation.slides)}

    def reorder_slides(self, presentation: Presentation, new_order: list[int]) -> dict:
        if sorted(new_order) != list(range(len(presentation.slides))):
            return {"success": False, "error": "Urutan tidak valid"}
        presentation.slides = [presentation.slides[i] for i in new_order]
        return {"success": True, "order": new_order}

    def duplicate_slide(self, presentation: Presentation, index: int) -> dict:
        if index < 0 or index >= len(presentation.slides):
            return {"success": False, "error": f"Index {index} di luar jangkauan"}
        original = presentation.slides[index]
        new_slide = Slide(
            title=original.title + " (copy)",
            content=original.content,
            layout=original.layout,
            notes=original.notes,
        )
        new_slide.images = original.images.copy()
        presentation.slides.insert(index + 1, new_slide)
        return {"success": True, "new_index": index + 1, "slide": new_slide.to_dict()}

    def create_from_outline(self, title: str, outline: list[dict],
                            author: str = "", theme: str = "modern") -> dict:
        pres = self.create_presentation(title=title, author=author, theme=theme)
        for item in outline:
            slide_title = item.get("title", "")
            slide_content = item.get("content", "")
            slide_layout = item.get("layout", "title_content")
            slide_notes = item.get("notes", "")
            slide = pres.add_slide(slide_title, slide_content, slide_layout, slide_notes)
            if "images" in item:
                slide.images = item["images"]

        return {"success": True, "presentation": pres.to_dict()}

    def get_presentation(self, title: str) -> Optional[dict]:
        for p in self.presentations:
            if p.title == title:
                return p.to_dict()
        return None

    def import_from_json(self, json_path: str) -> dict:
        try:
            with open(json_path, "r", encoding="utf-8") as f:
                data = json.load(f)

            pres = self.create_presentation(
                title=data.get("title", "Imported"),
                author=data.get("author", ""),
                theme=data.get("theme", "modern"),
            )

            for slide_data in data.get("slides", []):
                slide = pres.add_slide(
                    title=slide_data.get("title", ""),
                    content=slide_data.get("content", ""),
                    layout=slide_data.get("layout", "title_content"),
                    notes=slide_data.get("notes", ""),
                )
                slide.images = slide_data.get("images", [])

            return {"success": True, "presentation": pres.to_dict()}
        except Exception as e:
            return {"success": False, "error": str(e)}
